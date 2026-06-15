from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def set_premium_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(30, 41, 59) # Dark slate slate-800

def format_title(title_shape, text):
    title_shape.text = text
    title_frame = title_shape.text_frame
    for paragraph in title_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.name = 'Arial'
            run.font.size = Pt(36)
            run.font.bold = True
            run.font.color.rgb = RGBColor(56, 189, 248) # Sky-400

def format_content(content_shape, bullet_points):
    text_frame = content_shape.text_frame
    text_frame.clear()
    for point in bullet_points:
        p = text_frame.add_paragraph()
        p.text = point
        p.level = 0
        p.space_after = Pt(20)
        for run in p.runs:
            run.font.name = 'Arial'
            run.font.size = Pt(24)
            run.font.color.rgb = RGBColor(241, 245, 249) # Slate-100

prs = Presentation()

# Cover Slide
slide_layout = prs.slide_layouts[0] 
slide = prs.slides.add_slide(slide_layout)
set_premium_background(slide)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "PERANCANGAN BACKEND, DATABASE, DAN DASHBOARD ADMIN"
title_para = title.text_frame.paragraphs[0]
for run in title_para.runs:
    run.font.name = 'Arial'
    run.font.bold = True
    run.font.size = Pt(36)
    run.font.color.rgb = RGBColor(56, 189, 248) # Sky-400

subtitle.text = "Sistem Manajemen Produksi Event FND Production\n\nLaporan Kerja Praktek\nEdisyah Putra Waruwu (411231179)\nProgram Studi Teknik Informatika\nUniversitas Dian Nusantara - 2026"
for para in subtitle.text_frame.paragraphs:
    for run in para.runs:
        run.font.name = 'Arial'
        run.font.size = Pt(22)
        run.font.color.rgb = RGBColor(226, 232, 240)

slides_data = [
    {
        "title": "Latar Belakang",
        "content": [
            "Industri event management modern bergerak secara cepat dan dinamis.",
            "PT Fortuna Nusa Dream (FND Production) masih mengelola operasional secara manual dan parsial.",
            "Risiko administratif: overbooking alat, tumpang tindih jadwal penyewaan dan kru lapangan.",
            "Solusi: Pengembangan sistem informasi terintegrasi dengan backend modern (Next.js) dan database relasional (MySQL)."
        ]
    },
    {
        "title": "Tujuan Kerja Praktek",
        "content": [
            "Merancang dan membangun arsitektur basis data relasional (MySQL) yang optimal untuk operasional FND Production.",
            "Membangun dashboard admin berbasis web (Next.js) untuk mempermudah manajemen internal.",
            "Mengimplementasikan RESTful API routes yang aman sebagai jembatan komunikasi data sistem."
        ]
    },
    {
        "title": "Gambaran Umum Perusahaan",
        "content": [
            "Berdiri sejak 2015, FND Production adalah penyedia tata cahaya (lighting system) profesional.",
            "Melayani konser musik, wedding, corporate gathering, dan hiburan lainnya.",
            "Penempatan Kerja Praktek: Divisi Operasional dan Produksi (Technical Department).",
            "Peran: Tim Instalasi Lapangan & Analis Sistem untuk mengobservasi alur kerja dan merancang sistem manajemen."
        ]
    },
    {
        "title": "Arsitektur Sistem (Client-Server)",
        "content": [
            "Pendekatan Headless (berbasis API) untuk menjembatani Web Admin dan Mobile App (Klien & Kru).",
            "Frontend Layer: App Router Next.js, shadcn/ui, Tailwind CSS.",
            "API Layer: Next.js API Routes (Serverless) untuk logika bisnis dan transaksi database.",
            "Backend Layer: MySQL Database untuk pengolahan data relasional transaksional yang andal."
        ]
    },
    {
        "title": "Teknologi yang Digunakan",
        "content": [
            "Bahasa: TypeScript / JavaScript, SQL",
            "Framework: Next.js 16.2, React.js, Tailwind CSS 4.2",
            "Database: MySQL, mysql2/promise ORM",
            "Library Pendukung: shadcn/ui, React Hook Form, Zod, Recharts, Lucide React",
            "Infrastruktur: Vercel Cloud Deployment"
        ]
    },
    {
        "title": "Implementasi & Hasil Utama",
        "content": [
            "Dashboard Admin terpusat untuk kelola inventaris alat, kru, dan transaksi event secara real-time.",
            "Algoritma pencegahan konflik jadwal kru dan pencegahan overbooking peralatan.",
            "Database Transaction digunakan untuk memastikan integritas data apabila terjadi kegagalan sistem.",
            "Automasi update status inventaris saat terjadi perubahan pada jadwal (misal: event dibatalkan)."
        ]
    },
    {
        "title": "Evaluasi & Kesimpulan",
        "content": [
            "Sistem sukses mentransformasi proses persetujuan pesanan dari hitungan jam menjadi hitungan menit.",
            "Kelebihan: Antarmuka UI/UX modern, keamanan tipe data (TypeScript), database tangguh.",
            "Digitalisasi meminimalkan risiko human error dalam alokasi inventaris peralatan dan jadwal kru.",
            "Penerapan Next.js dan MySQL memberikan landasan performa dan integritas yang sangat solid."
        ]
    }
]

for s_data in slides_data:
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    set_premium_background(slide)
    
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    format_title(title_shape, s_data["title"])
    format_content(content_shape, s_data["content"])

prs.save("Presentasi_Laporan_KP_Premium.pptx")
print("Presentation generated successfully!")
